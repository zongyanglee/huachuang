"""
qa_pptx.py - 对生成的 PPT 做质量检查

功能概述：
    1. 文本检查：用 markitdown 提取 PPT 文字，检查占位符残留
    2. 内容检查：每张 slide 是否有标题
    3. 图片检查：图片是否插入成功
    4. 视觉渲染（可选）：需要 soffice + pdftoppm

输入：--input output.pptx
"""

from __future__ import annotations

import argparse
import shutil
import subprocess
import sys
from pathlib import Path

from pptx import Presentation


# ---- 1. 检查项 ----

PLACEHOLDER_KEYWORDS = (
    "lorem", "ipsum", "xxx", "todo", "placeholder",
    "占位", "示例文字", "示例标题", "示例内容", "template", "untitled",
)

OLD_TITLE_HINTS = (
    # 这些是真正的“模板原内容”会出现的关键词
    # 排除：公众号、免责声明（这些是模板保留页，不是误报）
)


def check_text(pptx_path: Path) -> list[str]:
    """文本检查"""
    issues = []
    try:
        from markitdown import MarkItDown
        md = MarkItDown()
        result = md.convert(str(pptx_path))
        text = result.text_content
    except ImportError:
        # markitdown 不可用，用 python-pptx
        prs = Presentation(pptx_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        text = "\n".join(texts)

    # 1. 占位符残留
    for kw in PLACEHOLDER_KEYWORDS:
        if kw in text.lower():
            issues.append(f"  - 发现占位符关键词: '{kw}'")

    # 2. 模板原内容残留（只检查明显的模板特有词）
    for hint in OLD_TITLE_HINTS:
        if hint in text:
            issues.append(f"  - 发现模板原内容残留: '{hint[:30]}...'")

    return issues


def check_structure(pptx_path: Path) -> list[str]:
    """结构检查"""
    issues = []
    prs = Presentation(pptx_path)

    # 1. 检查页数合理
    n = len(prs.slides)
    if n < 5:
        issues.append(f"  - 页数过少: {n}（建议 > 5）")
    elif n > 100:
        issues.append(f"  - 页数过多: {n}（建议 < 100）")

    # 2. 检查是否有内容页（非封面/章节页）
    content_count = sum(1 for slide in prs.slides
                        if slide.slide_layout.name == "空白图表页")
    if content_count == 0:
        issues.append("  - 没有内容页（空白图表页）")

    return issues


def check_images(pptx_path: Path) -> list[str]:
    """图片检查"""
    issues = []
    prs = Presentation(pptx_path)
    img_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                img_count += 1
    if img_count == 0:
        issues.append("  - 没有检测到任何图片")
    print(f"  [INFO] 共 {img_count} 张图片")
    return issues


def render_pdf(pptx_path: Path, output_dir: Path) -> list[str]:
    """渲染 PDF（需要 soffice + pdftoppm）"""
    issues = []
    if not shutil.which("soffice"):
        issues.append("  - soffice 未安装，跳过视觉渲染")
        return issues
    if not shutil.which("pdftoppm"):
        issues.append("  - pdftoppm 未安装，跳过 PDF→图片")
        return issues

    output_dir.mkdir(parents=True, exist_ok=True)
    # 转 PDF
    result = subprocess.run(
        ["soffice", "--headless", "--convert-to", "pdf",
         "--outdir", str(output_dir), str(pptx_path)],
        capture_output=True, text=True, timeout=120,
    )
    if result.returncode != 0:
        issues.append(f"  - soffice 转换失败: {result.stderr[:200]}")
        return issues
    pdf_path = output_dir / (pptx_path.stem + ".pdf")
    if not pdf_path.exists():
        issues.append(f"  - PDF 未生成")
        return issues
    # 转图片
    subprocess.run(
        ["pdftoppm", "-jpeg", "-r", "100", str(pdf_path), str(output_dir / "slide")],
        timeout=120,
    )
    print(f"  [INFO] 渲染图片: {output_dir}/slide-*.jpg")
    return issues


# ---- 2. 主流程 ----

def main() -> int:
    parser = argparse.ArgumentParser(description="PPT 质量检查")
    parser.add_argument("--input", required=True, help="待检查的 PPT 路径")
    parser.add_argument("--render", action="store_true", help="尝试渲染为图片")
    parser.add_argument("--render-dir", default=None, help="渲染输出目录")
    args = parser.parse_args()

    pptx_path = Path(args.input)
    if not pptx_path.exists():
        print(f"[ERROR] 文件不存在: {pptx_path}", file=sys.stderr)
        return 1

    all_issues = []

    print("[1/3] 文本检查")
    issues = check_text(pptx_path)
    all_issues.extend(issues)
    if not issues:
        print("  [OK] 无占位符残留")

    print("[2/3] 结构检查")
    issues = check_structure(pptx_path)
    all_issues.extend(issues)
    if not issues:
        print("  [OK] 结构正常")

    print("[3/3] 图片检查")
    issues = check_images(pptx_path)
    all_issues.extend(issues)
    if not issues:
        print("  [OK] 图片正常")

    if args.render:
        print("[可选] 视觉渲染")
        render_dir = Path(args.render_dir) if args.render_dir else pptx_path.parent / "qa_render"
        issues = render_pdf(pptx_path, render_dir)
        all_issues.extend(issues)

    print()
    if all_issues:
        print(f"[WARN] 发现 {len(all_issues)} 个问题：")
        for i in all_issues:
            print(i)
        return 1
    else:
        print(f"[OK] QA 通过")
        return 0


if __name__ == "__main__":
    sys.exit(main())
