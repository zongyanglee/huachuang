"""
build_pptx.py - 从 content.json + PPT 模板生成完整 PPT

完整 6 阶段管线:
1. 提取封面标题 + 解析内容(含加粗 runs)
2. 智能分页(图表 1:1 映射、风险提示合并、长文本拆分)
3. 构建纯文字 PPT(XML 层面编辑,保留模板格式)
4. 动态估算文本高度 + 插入图表标题 + 图片
5. QA(删除空图表标题)
6. 打包

用法:
    python build_pptx.py --template PPT模版.pptx --content _extracted/content.json --output output.pptx
"""

from __future__ import annotations

import copy
import html
import json
import re
import shutil
import sys
import tempfile
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET


# ── 命名空间 ──
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}
A = f'{{{NS["a"]}}}'
R = f'{{{NS["r"]}}}'
P = f'{{{NS["p"]}}}'


def inch(v: float) -> int:
    return int(v * 914400)


def estimate_text_height(text_lines: list, box_width_in: float) -> float:
    """Estimate the rendered body-text height in inches."""
    if not text_lines:
        return 0.3

    plain_lines = []
    for line in text_lines:
        if isinstance(line, str):
            plain_lines.append(line)
        elif isinstance(line, list):
            plain_lines.append("".join(r.get("text", "") for r in line))

    total_chars = sum(len(line) for line in plain_lines)
    n_paras = len(plain_lines)
    t = min(1.0, total_chars / 500)
    chars_per_inch = 5.5 - t * 1.0
    line_height = 0.23 + t * 0.07
    padding = 0.12 + t * 0.13
    para_spacing = max(0, n_paras - 1) * 0.06
    chars_per_line = max(1, int((box_width_in - 0.3) * chars_per_inch))
    total_lines = sum(max(1, -(-len(line) // chars_per_line)) for line in plain_lines)
    # PowerPoint's actual Chinese line metrics are often taller than the OOXML
    # nominal font size, especially with mixed bold runs. Keep a safety margin
    # so centered template text does not spill outside the estimated box.
    return max(0.3, (total_lines * line_height + padding + para_spacing) * 1.60)


# ── XML 工具 ──

def get_text(sp: ET.Element) -> str:
    return "".join(t.text or "" for t in sp.iter(f"{A}t")).strip()


def get_name(sp: ET.Element) -> str:
    cNvPr = sp.find(f".//{P}nvSpPr/{P}cNvPr")
    return cNvPr.get("name", "") if cNvPr is not None else ""


def set_text(sp: ET.Element, text: str) -> None:
    t_elems = list(sp.iter(f"{A}t"))
    if not t_elems:
        return
    t_elems[0].text = html.unescape(text)
    t_elems[0].set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
    for t in t_elems[1:]:
        t.text = ""


def set_multiline(sp: ET.Element, lines: list) -> None:
    """
    设置形状多行文本。
    lines 可以是 list[str](纯文本)或 list[list[dict]](runs 格式,
    每个 dict: {"text": str, "bold": bool})。
    """
    txBody = sp.find(f".//{P}txBody")
    if txBody is None:
        txBody = sp.find(f".//{A}txBody")
    if txBody is None:
        return
    paras = txBody.findall(f"{A}p")
    if not paras:
        return
    tpl_pPr = paras[0].find(f"{A}pPr")
    tpl_rPr = None
    tpl_r = paras[0].find(f"{A}r")
    if tpl_r is not None:
        tpl_rPr = tpl_r.find(f"{A}rPr")
    for p in paras[1:]:
        txBody.remove(p)
    first = paras[0]
    for r in list(first.findall(f"{A}r")):
        first.remove(r)
    for epr in list(first.findall(f"{A}endParaRPr")):
        first.remove(epr)
    if not lines:
        lines = [""]
    for i, line in enumerate(lines):
        if i == 0:
            p = first
        else:
            p = ET.SubElement(txBody, f"{A}p")
            if tpl_pPr is not None:
                p.append(copy.deepcopy(tpl_pPr))
        if isinstance(line, str):
            # 纯文本:单个 run
            r = ET.SubElement(p, f"{A}r")
            if tpl_rPr is not None:
                r.append(copy.deepcopy(tpl_rPr))
            t = ET.SubElement(r, f"{A}t")
            t.text = html.unescape(line)
            t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
        elif isinstance(line, list):
            # runs 格式:每个元素一个 run,支持加粗
            for run_data in line:
                r = ET.SubElement(p, f"{A}r")
                cloned = copy.deepcopy(tpl_rPr) if tpl_rPr is not None else ET.Element(f"{A}rPr")
                if run_data.get("bold"):
                    cloned.set("b", "1")
                r.append(cloned)
                t = ET.SubElement(r, f"{A}t")
                t.text = html.unescape(run_data.get("text", ""))
                t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")


def set_pos(sp: ET.Element, x: int, y: int, w: int, h: int) -> None:
    xfrm = sp.find(f".//{A}xfrm")
    if xfrm is None:
        return
    off = xfrm.find(f"{A}off")
    ext = xfrm.find(f"{A}ext")
    if off is not None:
        off.set("x", str(x))
        off.set("y", str(y))
    if ext is not None:
        ext.set("cx", str(w))
        ext.set("cy", str(h))


def set_pos_only(sp: ET.Element, x: int, y: int, w: int) -> None:
    xfrm = sp.find(f".//{A}xfrm")
    if xfrm is None:
        return
    off = xfrm.find(f"{A}off")
    ext = xfrm.find(f"{A}ext")
    if off is not None:
        off.set("x", str(x))
        off.set("y", str(y))
    if ext is not None:
        ext.set("cx", str(w))


def remove_pic(root: ET.Element) -> None:
    spTree = root.find(f".//{P}spTree")
    if spTree is None:
        return
    for pic in list(root.iter(f"{P}pic")):
        spTree.remove(pic)


def clone_spElement(sp: ET.Element) -> ET.Element:
    """深拷贝一个形状元素(保留所有子元素和属性)"""
    return copy.deepcopy(sp)


# ── 幻灯片分类与操作 ──

def classify_slide(root: ET.Element) -> str:
    sp_list = list(root.iter(f"{P}sp"))
    names = {get_name(sp) for sp in sp_list}
    texts = [get_text(sp) for sp in sp_list]
    joined = " ".join(texts)
    if "免责声明" in joined or "本材料仅供" in joined:
        return "closing"
    if "欢迎关注" in joined or "公众号" in joined:
        return "qrcode"
    if any("SECTION" in t for t in texts):
        return "section"
    if "周冠南" in joined or "SAC" in joined or "文本占位符 6" in names:
        return "cover"
    if "矩形 8" in names and "标题 1" in names:
        return "content"
    return "other"


def dup_slide(unpacked: Path, src_num: int) -> int:
    sd = unpacked / "ppt" / "slides"
    rd = sd / "_rels"
    src = sd / f"slide{src_num}.xml"
    src_r = rd / f"slide{src_num}.xml.rels"
    nums = [int(m.group(1)) for f in sd.glob("slide*.xml")
            if (m := re.match(r"slide(\d+)\.xml", f.name))]
    nxt = max(nums) + 1 if nums else 1
    shutil.copy2(src, sd / f"slide{nxt}.xml")
    if src_r.exists():
        shutil.copy2(src_r, rd / f"slide{nxt}.xml.rels")
        c = (rd / f"slide{nxt}.xml.rels").read_text(encoding="utf-8")
        c = re.sub(r'<[^>]*Relationship[^>]*Type="[^"]*notesSlide"[^>]*/>\s*', '', c)
        (rd / f"slide{nxt}.xml.rels").write_text(c, encoding="utf-8")
    return nxt


def reg_slide(unpacked: Path, num: int) -> None:
    pr = unpacked / "ppt" / "_rels" / "presentation.xml.rels"
    tree = ET.parse(pr)
    root = tree.getroot()
    max_r = max((int(m.group(1)) for r in root
                 if (m := re.match(r"rId(\d+)", r.get("Id", "")))), default=0)
    rid = f"rId{max_r + 1}"
    nr = ET.SubElement(root, f'{{{NS["rel"]}}}Relationship')
    nr.set("Id", rid)
    nr.set("Type", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide")
    nr.set("Target", f"slides/slide{num}.xml")
    tree.write(pr, encoding="utf-8", xml_declaration=True)

    pp = unpacked / "ppt" / "presentation.xml"
    tree = ET.parse(pp)
    root = tree.getroot()
    sldIdLst = root.find(f"{{{NS['p']}}}sldIdLst")
    if sldIdLst is None:
        sldIdLst = ET.SubElement(root, f"{{{NS['p']}}}sldIdLst")
    max_id = max((int(s.get("id", "0")) for s in sldIdLst), default=256)
    si = ET.SubElement(sldIdLst, f"{{{NS['p']}}}sldId")
    si.set("id", str(max_id + 1))
    si.set(f'{{{NS["r"]}}}id', rid)
    tree.write(pp, encoding="utf-8", xml_declaration=True)

    ct = unpacked / "[Content_Types].xml"
    tree = ET.parse(ct)
    root = tree.getroot()
    pn = f"/ppt/slides/slide{num}.xml"
    CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
    if not any(o.get("PartName") == pn for o in root.findall(f"{{{CT_NS}}}Override")):
        ov = ET.SubElement(root, f"{{{CT_NS}}}Override")
        ov.set("PartName", pn)
        ov.set("ContentType",
               "application/vnd.openxmlformats-officedocument.presentationml.slide+xml")
    tree.write(ct, encoding="utf-8", xml_declaration=True)


def del_slide(unpacked: Path, num: int) -> None:
    pp = unpacked / "ppt" / "presentation.xml"
    tree = ET.parse(pp)
    root = tree.getroot()
    sldIdLst = root.find(f"{{{NS['p']}}}sldIdLst")
    if sldIdLst is None:
        return
    pr = unpacked / "ppt" / "_rels" / "presentation.xml.rels"
    rt = ET.parse(pr)
    target_rid = None
    for r in rt.getroot():
        if f"slide{num}.xml" in r.get("Target", ""):
            target_rid = r.get("Id")
            break
    if target_rid:
        for si in list(sldIdLst):
            if si.get(f'{{{NS["r"]}}}id', "") == target_rid:
                sldIdLst.remove(si)
                break
    tree.write(pp, encoding="utf-8", xml_declaration=True)


# ── 封面标题提取 + 图表图片映射 ──

def generate_chart_image_map(items: list[dict]) -> dict:
    """从 content.json items 生成 {图表标题: 图片路径} 映射,按出现顺序一一对应。"""
    import re as _re
    chart_to_image = {}
    pending_charts, pending_images = [], []
    for item in items:
        if item["type"] == "chart_title":
            clean = _re.sub(r"^图表\s*\d+\s*", "", item["text"]).strip()
            pending_charts.append(clean)
        elif item["type"] == "image":
            pending_images.append(item["file"])
        else:
            for ci, ch in enumerate(pending_charts):
                if ci < len(pending_images) and ch not in chart_to_image:
                    chart_to_image[ch] = pending_images[ci]
            pending_charts, pending_images = [], []
    for ci, ch in enumerate(pending_charts):
        if ci < len(pending_images) and ch not in chart_to_image:
            chart_to_image[ch] = pending_images[ci]
    return chart_to_image


def extract_cover_title(docx_path: Path) -> str:
    """
    从原始 docx 的 XML 中提取报告标题
    查找前几个无样式段落中的标题文字
    """
    import re as re_m
    with zipfile.ZipFile(docx_path) as zf:
        xml = zf.read("word/document.xml").decode("utf-8")

    para_pattern = re_m.compile(r"<w:p[ >].*?</w:p>", re_m.DOTALL)
    paras_raw = para_pattern.findall(xml)

    title_parts = []
    for para_xml in paras_raw[:15]:
        text = "".join(re_m.findall(r'<w:t[^>]*>([^<]*)</w:t>', para_xml)).strip()
        style_match = re_m.search(r'<w:pStyle w:val="([^"]+)"/>', para_xml)
        style = style_match.group(1) if style_match else ""
        if not style and text:
            if text.startswith("【"):
                continue
            if text.startswith("--"):
                title_parts.append(text)
            elif len(text) >= 5:
                title_parts.append(text)

    if title_parts:
        return "\n".join(title_parts)
    return ""


# ── 内容解析 + 分页 ──

def parse_and_plan(items: list[dict], img_map: dict, cover_title: str) -> list[dict]:
    """
    从 content.json items 直接规划完整页面列表。
    
    分页规则（以图表为分隔）：
    - 每个图表（或连续 1-2 个图表）界定一页
    - 该页正文 = 当前图表之前、上个图表之后的所有 body 段落
    - 连续 > 2 个图表：前 2 个一页（附正文），其余独立成页（复用正文）
    - 无图表的正文：按 2 段/页分组，纯文页
    - 风险提示：独立页
    """
    pages = []
    pages.append({"type": "cover", "title": cover_title})
    # 长篇手册的三级标题较多。若全部提升为页标题，会造成大量只有一段
    # 正文的碎片页；此时保留二级标题作为页标题，并将三级标题嵌入正文
    # 首句且加粗，以同时保留层级和合理页数。
    dense_handbook = sum(
        1 for item in items
        if item.get("type") == "heading" and item.get("level", 1) >= 3
    ) > 20
    
    # 第一步：按 L1 标题分块，同时记录原文顺序用于页码
    chapters = []  # [(L1_heading, section_items)]
    current_chapter = None
    current_items = []
    
    for item in items:
        if item["type"] == "heading" and item.get("level") == 1:
            if current_chapter is not None:
                chapters.append((current_chapter, current_items))
            current_chapter = item["text"]
            current_items = []
        else:
            current_items.append(item)
    if current_chapter is not None:
        chapters.append((current_chapter, current_items))
    
    # 第二步：逐章分页
    risk_chapter_idx = None
    for ci, (ch_title, ch_items) in enumerate(chapters):
        if "风险提示" in ch_title:
            risk_chapter_idx = ci
            continue  # 风险提示后处理
        if "复盘" in ch_title:
            continue  # 跳过复盘章
        
        pages.append({"type": "section", "section_num": len([p for p in pages if p["type"]=="section"])+1,
                       "title": ch_title})
        
        # 收集本章的 body + 图表（图表属于其前面的 body）
        stream = []
        current_body = None
        sub_heading = ch_title  # 默认用章标题，遇到 L2 后切换
        pending_subheading = None
        
        for item in ch_items:
            if item["type"] == "heading" and item.get("level", 1) > 1:
                level = item.get("level", 1)
                if dense_handbook and level >= 3:
                    pending_subheading = item["text"]
                else:
                    sub_heading = item["text"]
                    pending_subheading = None
                continue
            elif item["type"] == "body":
                text = item["text"]
                if re.match(r"^注[：:]", text):
                    continue
                runs = item.get("runs", [])
                if pending_subheading:
                    prefix = f"{pending_subheading}："
                    text = prefix + text
                    runs = [{"text": prefix, "bold": True, "italic": False}] + runs
                    pending_subheading = None
                current_body = {"text": text, "runs": runs, "charts": [],
                                "title": sub_heading}  # 携带当前 L2/L3 标题
                stream.append(("body", current_body))
            elif item["type"] == "chart_title":
                clean = re.sub(r"^图表\s*\d+\s*", "", item["text"]).strip()
                if clean and current_body is not None:
                    current_body["charts"].append(clean)
        
        # 按图表边界分页
        i = 0
        while i < len(stream):
            body_parts = []
            chart_titles = []
            page_title = ch_title  # 默认章标题
            body_limit_reached = False
            body_chars = 0
            
            while i < len(stream):
                tp, data = stream[i]
                if tp == "body":
                    incoming_title = data.get("title") or ch_title
                    incoming_text = data["text"]
                    # 长篇手册常在两组图表之间包含多个小节。遇到标题变化、
                    # 段落过多或文字过长时先生成纯文字页，避免把数个小节和
                    # 后续图表强行挤到同一页并造成底部溢出。
                    if body_parts and (
                        incoming_title != page_title
                        or len(body_parts) >= 5
                        or body_chars + len(incoming_text) > 550
                    ):
                        body_limit_reached = True
                        break
                    if not body_parts:
                        page_title = incoming_title
                    body_parts.append((data["text"], data.get("runs", [])))
                    body_chars += len(incoming_text)
                    charts = data.get("charts", [])
                    chart_titles.extend([c for c in charts if c.strip()])
                    if charts:
                        i += 1
                        break
                i += 1

            if body_limit_reached:
                pages.append({
                    "type": "content", "title": page_title,
                    "body": [t for t, r in body_parts],
                    "body_runs": [r for t, r in body_parts],
                    "charts": [],
                })
                continue
            
            if not chart_titles and i >= len(stream):
                # 剩余 body — 纯文页
                while i < len(stream):
                    tp, data = stream[i]
                    if tp == "body":
                        body_parts.append((data["text"], data.get("runs", [])))
                    i += 1
                for bi in range(0, len(body_parts), 2):
                    chunk = body_parts[bi:bi+2]
                    pages.append({
                        "type": "content", "title": page_title,
                        "body": [t for t, r in chunk],
                        "body_runs": [r for t, r in chunk],
                        "charts": [],
                    })
                break
            
            if not chart_titles:
                continue

            # 图表页的可用正文高度显著小于纯文字页。若同一小节在图表前有
            # 多段说明，先把较早的段落放到纯文字页，仅保留最贴近图表的
            # 1–2 段说明，避免图表被推到页外或缩成不可读的细条。
            chart_body = list(body_parts)
            leading_body = []
            while len(chart_body) > 1 and (
                len(chart_body) > 2
                or sum(len(t) for t, _ in chart_body) > 450
            ):
                leading_body.append(chart_body.pop(0))
            if leading_body:
                pages.append({
                    "type": "content", "title": page_title,
                    "body": [t for t, r in leading_body],
                    "body_runs": [r for t, r in leading_body],
                    "charts": [],
                })
            
            # 图表分页
            for ci_offset in range(0, len(chart_titles), 2):
                batch = chart_titles[ci_offset:ci_offset+2]
                # Word 中偶尔会把普通规则小标题误标成“图表标题”。只有在
                # chart_image_map 中确实有对应图片时，才创建图表模块，避免
                # 生成只有蓝色标题条、没有图表内容的空框。
                chart_entries = [
                    {"title": ct, "images": [img_map[ct]]}
                    for ct in batch if img_map.get(ct)
                ]
                pages.append({
                    "type": "content", "title": page_title,
                    "body": [t for t, r in chart_body],
                    "body_runs": [r for t, r in chart_body],
                    "charts": chart_entries,
                })
    
    # ── 风险提示（硬性规则：下周关注 + 流动性警告 合并为一页，标题固定"风险提示"）──
    if risk_chapter_idx is not None:
        _, risk_items = chapters[risk_chapter_idx]
        risk_body = []
        risk_runs = []
        for item in risk_items:
            if item["type"] == "body" and not re.match(r"^注[：:]|^图表", item["text"]):
                risk_body.append(item["text"])
                risk_runs.append(item.get("runs", []))
        
        # 从上一章末尾提取"下周关注"并入风险提示页
        for pi in range(len(pages) - 1, -1, -1):
            if pages[pi]["type"] == "content" and pages[pi].get("body"):
                last_body = pages[pi]["body"][-1]
                last_runs = pages[pi]["body_runs"][-1]
                if last_body.strip().startswith("下周关注"):
                    # 从原页移除
                    pages[pi]["body"] = pages[pi]["body"][:-1]
                    pages[pi]["body_runs"] = pages[pi]["body_runs"][:-1]
                    # 清理空页
                    if not pages[pi]["body"] and not pages[pi].get("charts"):
                        pages.pop(pi)
                    # 前置到风险提示
                    risk_body.insert(0, last_body)
                    risk_runs.insert(0, last_runs)
                break
        
        if risk_body:
            pages.append({
                "type": "content",
                "title": "风险提示",
                "body": risk_body,
                "body_runs": risk_runs,
                "charts": [],
            })
    
    pages.append({"type": "qrcode"})
    pages.append({"type": "closing"})
    return pages

# ── 图片插入辅助 ──

def _insert_image(unpacked_dir: Path, root: ET.Element, snum: int,
                  ppt_media: Path, rels: dict, img_src: str,
                  x_emu: int, y_emu: int, w_emu: int, max_rid: int,
                  pic_id: int) -> str:
    """将图片复制到 PPT media 目录,添加到 slide XML,注册 rels 和 Content_Types。"""
    import shutil as _shutil
    src_path = Path(img_src)
    if not src_path.exists():
        return ""

    # 跳过无法嵌入的格式
    ext = src_path.suffix.lower()
    if ext == ".svg":
        print(f"  [SKIP] 不支持的图片格式: {src_path.name}")
        return ""

    # EMF/WMF 是 PowerPoint 原生格式,直接嵌入
    from PIL import Image as PILImage
    try:
        with PILImage.open(src_path) as pil_img:
            pw, ph = pil_img.size
            h_emu = int(w_emu * ph / pw) if pw > 0 else w_emu
    except Exception:
        h_emu = w_emu

    # 限制图片高度不超出页面(底部留 0.3" 边距)
    max_h_emu = int((7.5 - 0.3) * 914400) - y_emu
    if max_h_emu > 0 and h_emu > max_h_emu:
        h_emu = max_h_emu

    img_name = f"docx_{src_path.name}"
    dst = ppt_media / img_name
    if not dst.exists():
        _shutil.copy2(src_path, dst)

    # 注册 Content_Type
    ct_file = unpacked_dir / "[Content_Types].xml"
    ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    if ct_file.exists():
        ct_tree = ET.parse(ct_file)
        ct_root = ct_tree.getroot()
        part_name = f"/ppt/media/{img_name}"
        ext_map = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                   ".gif": "image/gif", ".bmp": "image/bmp",
                   ".emf": "image/x-emf", ".wmf": "image/x-wmf"}
        content_type = ext_map.get(ext, "image/png")
        if not any(o.get("PartName") == part_name for o in ct_root.findall(f"{{{ct_ns}}}Override")):
            ov = ET.SubElement(ct_root, f"{{{ct_ns}}}Override")
            ov.set("PartName", part_name)
            ov.set("ContentType", content_type)
        ct_tree.write(ct_file, encoding="utf-8", xml_declaration=True)

    rid = f"rId{max_rid + 1}"
    rels[rid] = f"../media/{img_name}"

    add_image_to_slide(root, img_name, x_emu, y_emu, w_emu, h_emu, rid, pic_id)
    return rid


def _save_rels(rels_path: Path, rels: dict[str, str]) -> None:
    """将新的图片关系写入 rels 文件(保留原有关系)"""
    if not rels_path.exists():
        return
    tree = ET.parse(rels_path)
    root = tree.getroot()
    existing_ids = {r.get("Id") for r in root}
    for rid, target in rels.items():
        if rid not in existing_ids:
            nr = ET.SubElement(root, f'{{{NS["rel"]}}}Relationship')
            nr.set("Id", rid)
            nr.set("Type", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image")
            nr.set("Target", target)
    tree.write(rels_path, encoding="utf-8", xml_declaration=True)


def add_image_to_slide(root: ET.Element, img_name: str, x_emu: int, y_emu: int,
                       w_emu: int, h_emu: int, rid: str, pic_id: int) -> None:
    """在 slide XML 的 spTree 中插入图片元素(pic_id 必须全局唯一)"""
    spTree = root.find(f".//{P}spTree")
    if spTree is None:
        return
    pic_str = (
        f'<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        f' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        f'<p:nvPicPr><p:cNvPr id="{pic_id}" name="{img_name}"/>'
        f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr><p:nvPr/></p:nvPicPr>'
        f'<p:blipFill><a:blip r:embed="{rid}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
        f'<p:spPr><a:xfrm><a:off x="{x_emu}" y="{y_emu}"/>'
        f'<a:ext cx="{w_emu}" cy="{h_emu}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>'
    )
    spTree.append(ET.fromstring(pic_str))


# ── 命名空间修复 ──

def fix_ns(unpacked: Path) -> None:
    ns_map = {
        'http://schemas.openxmlformats.org/presentationml/2006/main': 'p',
        'http://schemas.openxmlformats.org/drawingml/2006/main': 'a',
        'http://schemas.openxmlformats.org/officeDocument/2006/relationships': 'r',
        'http://schemas.openxmlformats.org/package/2006/content-types': '',
        'http://schemas.openxmlformats.org/package/2006/relationships': '',
    }
    for f in unpacked.rglob("*.xml"):
        c = f.read_text(encoding="utf-8")
        mod = False
        for uri, pfx in ns_map.items():
            m = re.search(f'xmlns:ns(\\d+)="{re.escape(uri)}"', c)
            if m:
                n = m.group(1)
                if pfx:
                    c = c.replace(f'xmlns:ns{n}=', f'xmlns:{pfx}=')
                    c = c.replace(f'<ns{n}:', f'<{pfx}:').replace(f'</ns{n}:', f'</{pfx}:')
                else:
                    c = c.replace(f'xmlns:ns{n}=', 'xmlns=')
                    c = c.replace(f'<ns{n}:', '<').replace(f'</ns{n}:', '</')
                mod = True
        c = re.sub(r'\bns(\d+):(id|embed|link)\b', r'r:\2', c)
        if mod or 'ns0:' in c or 'ns1:' in c:
            f.write_text(c, encoding="utf-8")

    for f in unpacked.rglob("*.rels"):
        c = f.read_text(encoding="utf-8")
        mod = False
        for uri, pfx in ns_map.items():
            m = re.search(f'xmlns:ns(\\d+)="{re.escape(uri)}"', c)
            if m:
                n = m.group(1)
                if pfx:
                    c = c.replace(f'xmlns:ns{n}=', f'xmlns:{pfx}=')
                    c = c.replace(f'<ns{n}:', f'<{pfx}:').replace(f'</ns{n}:', f'</{pfx}:')
                else:
                    c = c.replace(f'xmlns:ns{n}=', 'xmlns=')
                    c = c.replace(f'<ns{n}:', '<').replace(f'</ns{n}:', '</')
                mod = True
        c = re.sub(r'<(/?)\\s*ns\\d+:', r'<\\1', c)
        if mod or 'ns0:' in c or 'ns1:' in c:
            f.write_text(c, encoding="utf-8")


# ── 主函数 ──

def main() -> int:
    import argparse
    parser = argparse.ArgumentParser(description="从 content.json + PPT 模板生成完整 PPT")
    parser.add_argument("--template", required=True, help="PPT 模板文件路径")
    parser.add_argument("--content", required=True, help="content.json 所在目录(含 chart_image_map.json)")
    parser.add_argument("--docx", default=None, help="原始 docx 路径(用于提取封面标题,可选)")
    parser.add_argument("--output", default="output.pptx", help="输出 PPT 路径")
    args = parser.parse_args()

    template_path = Path(args.template).expanduser().resolve()
    content_dir = Path(args.content).expanduser().resolve()
    content_json = content_dir / "content.json"
    output_path = Path(args.output).expanduser().resolve()

    if not template_path.exists():
        print(f"[ERROR] 模板不存在: {template_path}", file=sys.stderr)
        return 1
    if not content_json.exists():
        print(f"[ERROR] content.json 不存在: {content_json}", file=sys.stderr)
        return 1

    output_path.parent.mkdir(parents=True, exist_ok=True)
    # Use the OS temporary directory. Some managed workspaces permit final
    # output writes but reject OOXML extraction files such as
    # ``[Content_Types].xml`` inside the project output directory.
    unpacked_dir = Path(tempfile.mkdtemp(prefix="docx_to_pptx_"))

    # ── Phase 1: 提取封面标题 + 解析内容 ──
    cover_title = ""
    if args.docx and Path(args.docx).exists():
        cover_title = extract_cover_title(Path(args.docx))
    else:
        # 无法提取封面标题时,尝试从第一个章节标题推导
        cover_title = "报告标题"  # fallback
    print(f"[OK] 封面标题:\n    {cover_title.replace(chr(10), chr(10)+'    ')}")

    with open(content_json, "r", encoding="utf-8") as f:
        items = json.load(f)

    # 生成图表→图片映射
    chart_img_map = generate_chart_image_map(items)
    print(f"[OK] 加载图表→图片映射: {len(chart_img_map)} 个")

    pages = parse_and_plan(items, chart_img_map, cover_title)
    print(f"[OK] 规划 {len(pages)} 页幻灯片")

    # ── Phase 2: 解压模板 + 构建纯文字 PPT ──
    with zipfile.ZipFile(template_path) as zf:
        zf.extractall(unpacked_dir)
    fix_ns(unpacked_dir)

    sd = unpacked_dir / "ppt" / "slides"
    rd = sd / "_rels"

    # 分析模板
    sfs = sorted(sd.glob("slide*.xml"),
                 key=lambda f: int(re.search(r"(\d+)", f.stem).group(1)))
    template_slides = []
    for sf in sfs:
        tree = ET.parse(sf)
        role = classify_slide(tree.getroot())
        num = int(re.search(r"(\d+)", sf.stem).group(1))
        template_slides.append({"file": sf, "num": num, "role": role})

    ref_cover = next((s for s in template_slides if s["role"] == "cover"), None)
    ref_section = next((s for s in template_slides if s["role"] == "section"), None)
    ref_content = next((s for s in template_slides if s["role"] == "content"), None)
    ref_qrcode = next((s for s in template_slides if s["role"] == "qrcode"), None)
    ref_closing = next((s for s in template_slides if s["role"] == "closing"), None)

    if not all([ref_cover, ref_section, ref_content, ref_qrcode, ref_closing]):
        print("[ERROR] 模板缺少必要的幻灯片类型", file=sys.stderr)
        return 1

    ref_map = {
        "cover": ref_cover["num"],
        "section": ref_section["num"],
        "content": ref_content["num"],
        "qrcode": ref_qrcode["num"],
        "closing": ref_closing["num"],
    }

    # 注册 EMF/WMF 支持(PowerPoint 原生格式)
    ct_file = unpacked_dir / "[Content_Types].xml"
    ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    ct_tree = ET.parse(ct_file)
    ct_root = ct_tree.getroot()
    for ext, mime in [("emf", "image/x-emf"), ("wmf", "image/x-wmf")]:
        if not any(d.get("Extension") == ext for d in ct_root.findall(f"{{{ct_ns}}}Default")):
            d = ET.SubElement(ct_root, f"{{{ct_ns}}}Default")
            d.set("Extension", ext)
            d.set("ContentType", mime)
    ct_tree.write(ct_file, encoding="utf-8", xml_declaration=True)

    # 复制幻灯片
    assignments = []
    content_indices = []  # 记录哪些 assignment 是正文页
    for idx, page in enumerate(pages):
        ref_num = ref_map[page["type"]]
        n = dup_slide(unpacked_dir, ref_num)
        reg_slide(unpacked_dir, n)
        assignments.append((n, page))
        if page["type"] == "content":
            content_indices.append(idx)

    # 删除模板原幻灯片
    keep_nums = {n for n, _ in assignments}
    for ts in template_slides:
        if ts["num"] not in keep_nums:
            del_slide(unpacked_dir, ts["num"])
            sf = sd / f"slide{ts['num']}.xml"
            sr = rd / f"slide{ts['num']}.xml.rels"
            if sf.exists():
                sf.unlink()
            if sr.exists():
                sr.unlink()

    # 填充文字
    margin_left = 403237  # 匹配模板"矩形 8" left
    content_width = inch(13.333) - margin_left - inch(0.5)

    for snum, page in assignments:
        sf = sd / f"slide{snum}.xml"
        if not sf.exists():
            continue
        tree = ET.parse(sf)
        root = tree.getroot()
        shapes = {get_name(sp): sp for sp in root.iter(f"{P}sp") if get_name(sp)}
        ptype = page["type"]

        if ptype == "cover":
            title_text = page.get("title", "")
            if "内容占位符 9" in shapes:
                # 使用多行模式,支持标题换行
                if "\n" in title_text:
                    set_multiline(shapes["内容占位符 9"], title_text.split("\n"))
                else:
                    set_text(shapes["内容占位符 9"], title_text)

        elif ptype == "section":
            sec_num = page.get("section_num", 1)
            sec_title = page.get("title", "")
            if "文本框 15" in shapes:
                set_text(shapes["文本框 15"], f"SECTION    {sec_num}")
            if "Rectangle 1" in shapes:
                set_text(shapes["Rectangle 1"], sec_title)

        elif ptype == "content":
            page_title = page.get("title", "")
            body = page.get("body", [])

            if "标题 1" in shapes:
                set_text(shapes["标题 1"], page_title)
                set_pos(shapes["标题 1"], margin_left, inch(0.3), content_width, inch(0.5))

            if "矩形 8" in shapes:
                # 优先用 runs 格式,否则纯文本;空正文时清空模板残留
                run_lines = page.get("body_runs", [])
                if run_lines and any(any(r.get("bold") for r in rl) for rl in run_lines if rl):
                    set_multiline(shapes["矩形 8"], run_lines)
                    height_lines = run_lines
                elif body:
                    set_multiline(shapes["矩形 8"], body)
                    height_lines = body
                else:
                    set_multiline(shapes["矩形 8"], [""])
                    height_lines = [""]
                body_pr = shapes["矩形 8"].find(f".//{A}bodyPr")
                if body_pr is not None:
                    body_pr.set("anchor", "t")
                    for tag in ("spAutoFit", "normAutofit", "noAutofit"):
                        for child in list(body_pr.findall(f"{A}{tag}")):
                            body_pr.remove(child)
                    ET.SubElement(body_pr, f"{A}noAutofit")
                body_height = min(
                    5.85,
                    estimate_text_height(height_lines, content_width / 914400),
                )
                set_pos(
                    shapes["矩形 8"], margin_left, 1043034,
                    content_width, inch(body_height),
                )

            # 清空图表标题(Phase 4 会重新填充)
            if "内容占位符 11" in shapes:
                set_text(shapes["内容占位符 11"], "")
                set_pos(shapes["内容占位符 11"], 0, inch(9), inch(0.1), inch(0.1))

            if "Rectangle 2" in shapes:
                set_text(shapes["Rectangle 2"], "")

            if "灯片编号占位符 2" in shapes:
                sp = shapes["灯片编号占位符 2"]
                xfrm = sp.find(f".//{A}xfrm")
                if xfrm is not None:
                    set_pos(shapes["灯片编号占位符 2"], 0, inch(9), inch(0.1), inch(0.1))

            remove_pic(root)

        tree.write(sf, encoding="utf-8", xml_declaration=True)

    print("[OK] Phase 2 完成:文字填充")

    # ── Phase 3: 打包临时 PPT + 读取文本框位置 ──
    fix_ns(unpacked_dir)
    with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for fp in sorted(unpacked_dir.rglob("*")):
            if fp.is_file():
                zf.write(fp, fp.relative_to(unpacked_dir))
    print("[OK] Phase 3 完成:临时 PPT 打包")

    # 收集每个正文页的文本框位置(估算版)
    from pptx import Presentation as PptxPresentation
    prs = PptxPresentation(str(output_path))

    body_bottoms = {}
    slide_idx = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name == '矩形 8':
                left_in = shape.left / 914400
                top_in = shape.top / 914400
                width_in = shape.width / 914400
                xml_height = shape.height / 914400  # XML 中的高度(不可靠)

                # 估算实际高度
                text_content = shape.text_frame.text
                lines = text_content.split('\n')
                est_height = estimate_text_height(lines, width_in)
                # 至少取 XML 高度
                height_in = max(xml_height, est_height)
                bottom_in = top_in + height_in

                body_bottoms[slide_idx] = {
                    "left": left_in,
                    "top": top_in,
                    "width": width_in,
                    "height": height_in,
                    "bottom": bottom_in,
                }
                break
        slide_idx += 1

    print(f"[OK] Phase 3 完成:估算 {len(body_bottoms)} 个文本框位置")

    # ── Phase 4: 插入图表标题 + 图片 ──
    # 布局参数(居中对称)
    SLIDE_W = 13.333
    MARGIN = 0.5
    CONTENT_W = SLIDE_W - 2 * MARGIN
    CHART_H = inch(0.35)
    CHART_W_SINGLE = inch(10.0)
    CHART_W_DUAL = inch(5.0)
    CHART_GAP_HORIZ = inch(0.3)
    MAX_CHARTS = 2
    BASE_GAP = 0.35
    ppt_media = unpacked_dir / "ppt" / "media"
    ppt_media.mkdir(exist_ok=True)

    global_pic_id = 100  # 全局递增,确保每张图片 ID 唯一

    for ci in content_indices:
        snum, page = assignments[ci]
        sf = sd / f"slide{snum}.xml"
        if not sf.exists():
            continue

        charts = page.get("charts", [])
        # 兼容新旧格式:旧格式是 [str],新格式是 [{"title": str, "images": [str]}]
        if charts and isinstance(charts[0], str):
            charts = [{"title": c, "images": []} for c in charts]

        charts = [c for c in charts if c.get("title", "").strip()]
        if not charts:
            continue

        seen = set()
        unique_charts = []
        for c in charts:
            if c["title"] not in seen:
                seen.add(c["title"])
                unique_charts.append(c)
        charts = unique_charts[:MAX_CHARTS]

        slide_position = ci
        body_info = body_bottoms.get(slide_position)
        if body_info is None:
            continue

        tree = ET.parse(sf)
        root = tree.getroot()
        shapes = {get_name(sp): sp for sp in root.iter(f"{P}sp") if get_name(sp)}

        if "内容占位符 11" in shapes:
            set_pos(shapes["内容占位符 11"], 0, inch(9), inch(0.1), inch(0.1))

        spTree = root.find(f".//{P}spTree")
        if spTree is None:
            continue

        tpl_chart_sp = shapes.get("内容占位符 11")
        if tpl_chart_sp is None:
            continue

        est_h = body_info["height"]
        dynamic_gap = inch(BASE_GAP + min(0.25, est_h * 0.12))
        chart_y = inch(body_info["bottom"]) + dynamic_gap
        n_charts = len(charts)

        # 加载 rels
        rels_path = rd / f"slide{snum}.xml.rels"
        rels = {}
        if rels_path.exists():
            rt = ET.parse(rels_path)
            rels = {r.get("Id"): r.get("Target", "") for r in rt.getroot()}

        # ── 第二步:在图表标题下方插入图片 ──
        max_rid = max((int(m.group(1)) for r in rels if (m := __import__('re').match(r"rId(\d+)", r))), default=0)
        title_info = []  # [(x, y_emu, w_emu, chart_entry)]

        if n_charts == 1:
            x = inch(MARGIN + (CONTENT_W - 10.0) / 2)
            new_sp = copy.deepcopy(tpl_chart_sp)
            set_text(new_sp, charts[0]["title"])
            set_pos(new_sp, x, chart_y, CHART_W_SINGLE, CHART_H)
            spTree.append(new_sp)
            title_info.append((x, chart_y, CHART_W_SINGLE, charts[0]))

        elif n_charts == 2:
            total_dual_w = 5.0 * 2 + 0.3
            left_start = MARGIN + (CONTENT_W - total_dual_w) / 2
            for c_idx, ch_entry in enumerate(charts):
                new_sp = copy.deepcopy(tpl_chart_sp)
                set_text(new_sp, ch_entry["title"])
                x = inch(left_start + c_idx * (5.0 + 0.3))
                set_pos(new_sp, x, chart_y, CHART_W_DUAL, CHART_H)
                spTree.append(new_sp)
                title_info.append((x, chart_y, CHART_W_DUAL, ch_entry))

        # ── 在图表标题下方插入图片 ──
        TITLE_IMG_GAP = inch(0.1)
        for tp_x, tp_y, tp_w, ch_entry in title_info:
            img_paths = ch_entry.get("images", [])
            valid_imgs = [p for p in img_paths if p and Path(p).exists()]
            if not valid_imgs:
                continue

            img_y = tp_y + CHART_H + TITLE_IMG_GAP
            n_imgs = len(valid_imgs)

            if n_imgs == 1:
                max_rid += 1; global_pic_id += 1
                _insert_image(unpacked_dir, root, snum, ppt_media, rels,
                             valid_imgs[0], tp_x, img_y, tp_w, max_rid - 1, global_pic_id - 1)
            else:
                sub_w = int((tp_w - CHART_GAP_HORIZ * (n_imgs - 1)) / n_imgs)
                for ii, img_path in enumerate(valid_imgs):
                    ix = tp_x + ii * (sub_w + CHART_GAP_HORIZ)
                    max_rid += 1; global_pic_id += 1
                    _insert_image(unpacked_dir, root, snum, ppt_media, rels,
                                 img_path, ix, img_y, sub_w, max_rid - 1, global_pic_id - 1)

        _save_rels(rels_path, rels)
        tree.write(sf, encoding="utf-8", xml_declaration=True)

    print(f"[OK] Phase 4 完成:图表标题 + 图片插入(动态间距)")

    # ── Phase 5: QA - 删除残留的空图表标题(在 XML 层面操作)──
    empty_removed = 0
    for snum, page in assignments:
        if page["type"] != "content":
            continue
        sf = sd / f"slide{snum}.xml"
        if not sf.exists():
            continue
        tree = ET.parse(sf)
        root = tree.getroot()
        spTree = root.find(f".//{P}spTree")
        if spTree is None:
            continue
        for sp in list(spTree.findall(f"{P}sp")):
            if get_name(sp) == "内容占位符 11":
                text = get_text(sp)
                if not text.strip():
                    spTree.remove(sp)
                    empty_removed += 1
        tree.write(sf, encoding="utf-8", xml_declaration=True)
    print(f"[QA] 删除了 {empty_removed} 个空图表标题模块" if empty_removed else f"[QA] 无空图表标题 ✓")

    # ── Phase 6: 最终打包 ──
    fix_ns(unpacked_dir)
    with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for fp in sorted(unpacked_dir.rglob("*")):
            if fp.is_file():
                zf.write(fp, fp.relative_to(unpacked_dir))

    # 统计
    cover_count = sum(1 for _, p in assignments if p["type"] == "cover")
    section_count = sum(1 for _, p in assignments if p["type"] == "section")
    content_count = sum(1 for _, p in assignments if p["type"] == "content")
    # 统计实际插入的图表标题(每页最多 2 个,去重后)
    chart_count = 0
    for _, p in assignments:
        if p["type"] == "content":
            charts = p.get("charts", [])
            if charts and isinstance(charts[0], str):
                charts = [{"title": c} for c in charts]
            chart_count += len([c for c in charts if c.get("title", "").strip()])
    chart_count = min(chart_count, content_count * 2)  # 理论上限

    print(f"[OK] 最终完成: {output_path}")
    print(f"     封面: {cover_count}, 章节分隔: {section_count}, 正文: {content_count}")
    print(f"     图表标题: {chart_count}, 总页数: {len(assignments)}")

    # 清理
    if unpacked_dir.exists():
        shutil.rmtree(unpacked_dir)

    return 0


if __name__ == "__main__":
    sys.exit(main())
